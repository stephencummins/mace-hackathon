"""Shared visual primitives for presentations in this folder.

One stylesheet means decks stay visually consistent. Each `build_*.py`
imports the constants and helpers it needs from here.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Visual identity ---------------------------------------------------------

NAVY = RGBColor(0x0B, 0x2E, 0x4F)          # primary background accent
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)      # body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xC8, 0x4E, 0x2C)        # terracotta accent
MUTED = RGBColor(0x6B, 0x6B, 0x6B)         # secondary text
GREEN = RGBColor(0x2E, 0x7D, 0x32)         # results / positive
LIGHT_BG = RGBColor(0xF5, 0xF3, 0xEE)      # panel background

FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


# --- Primitives --------------------------------------------------------------

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    bold=False,
    color=CHARCOAL,
    align=PP_ALIGN.LEFT,
    font=FONT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    title_size=20,
    body_size=14,
    color=CHARCOAL,
    accent=ACCENT,
):
    """Add a vertical list of bullets.

    Each item is either a plain string (single line) or a (headline, sub)
    tuple — headline bold, sub muted on the next line.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(8)

        if isinstance(item, tuple):
            head, sub = item
            p.add_run().text = "  "
            r2 = p.add_run()
            r2.text = head
            r2.font.name = FONT
            r2.font.size = Pt(title_size)
            r2.font.bold = True
            r2.font.color.rgb = color

            sub_p = tf.add_paragraph()
            sub_p.space_after = Pt(10)
            sub_run = sub_p.add_run()
            sub_run.text = "      " + sub
            sub_run.font.name = FONT
            sub_run.font.size = Pt(body_size)
            sub_run.font.color.rgb = MUTED
        else:
            run = p.add_run()
            run.text = "•  " + item
            run.font.name = FONT
            run.font.size = Pt(body_size)
            run.font.color.rgb = color


def header_band(slide, title, subtitle=None):
    """Top navy band with white title text plus an accent stripe."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), NAVY)
    add_rect(slide, 0, Inches(1.2), SLIDE_W, Inches(0.08), ACCENT)
    add_text(
        slide,
        Inches(0.5),
        Inches(0.25),
        SLIDE_W - Inches(1),
        Inches(0.7),
        title,
        size=32,
        bold=True,
        color=WHITE,
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.5),
            Inches(0.78),
            SLIDE_W - Inches(1),
            Inches(0.4),
            subtitle,
            size=16,
            color=WHITE,
        )


def footer_band(slide, slide_num, total=10, attribution="M+AI+CE Hackathon  ·  Mace"):
    add_text(
        slide,
        Inches(0.5),
        SLIDE_H - Inches(0.45),
        Inches(8),
        Inches(0.3),
        attribution,
        size=10,
        color=MUTED,
    )
    add_text(
        slide,
        SLIDE_W - Inches(2),
        SLIDE_H - Inches(0.45),
        Inches(1.5),
        Inches(0.3),
        f"{slide_num} / {total}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
